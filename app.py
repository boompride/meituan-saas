from flask import Flask, render_template, request, redirect, send_file, session, url_for
from werkzeug.utils import secure_filename
import os
import pandas as pd
from pptx import Presentation
from fpdf import FPDF
import uuid

app = Flask(__name__)
app.secret_key = 'demo-secret-key'

UPLOAD_FOLDER = 'uploads'
REPORT_FOLDER = 'reports'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(REPORT_FOLDER, exist_ok=True)

@app.route('/', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form['username']
        session['username'] = username
        return redirect('/dashboard')
    return render_template('login.html')

@app.route('/dashboard')
def dashboard():
    if 'username' not in session:
        return redirect('/')
    return render_template('dashboard.html', username=session['username'])

@app.route('/upload', methods=['POST'])
def upload():
    file = request.files['file']
    export_format = request.form.get('format', 'pptx')
    if file:
        filename = secure_filename(file.filename)
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        file.save(filepath)
        if export_format == 'pdf':
            report_path = generate_pdf_report(filepath)
        else:
            report_path = generate_pptx_report(filepath)
        return send_file(report_path, as_attachment=True)
    return '上传失败'

def generate_pptx_report(filepath):
    df = pd.read_excel(filepath)
    df.columns = df.iloc[0]
    df = df[1:].reset_index(drop=True)

    for col in ["总收入（元）", "平台技术服务费（元）", "商家营销费用（元）", "已消费后退款（元）"]:
        df[col] = pd.to_numeric(df[col], errors='coerce')

    total_income = df['总收入（元）'].sum()
    refund = df['已消费后退款（元）'].sum()
    platform_fee = df['平台技术服务费（元）'].sum()
    marketing = df['商家营销费用（元）'].sum()
    actual = total_income - refund - platform_fee - marketing

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "经营分析报告"
    slide.placeholders[1].text = f"总收入：¥{total_income:.2f}\n实际收入：¥{actual:.2f}"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "建议与分析"
    refund_rate = refund / total_income if total_income else 0
    notes = f"退款率：{refund_rate:.2%}\n"
    if refund_rate > 0.15:
        notes += "退款率偏高，建议优化套餐说明和服务流程"
    else:
        notes += "退款率正常"
    slide2.placeholders[1].text = notes

    out_path = os.path.join(REPORT_FOLDER, f"report_{uuid.uuid4().hex}.pptx")
    prs.save(out_path)
    return out_path

def generate_pdf_report(filepath):
    df = pd.read_excel(filepath)
    df.columns = df.iloc[0]
    df = df[1:].reset_index(drop=True)

    for col in ["总收入（元）", "平台技术服务费（元）", "商家营销费用（元）", "已消费后退款（元）"]:
        df[col] = pd.to_numeric(df[col], errors='coerce')

    total_income = df['总收入（元）'].sum()
    refund = df['已消费后退款（元）'].sum()
    platform_fee = df['平台技术服务费（元）'].sum()
    marketing = df['商家营销费用（元）'].sum()
    actual = total_income - refund - platform_fee - marketing
    refund_rate = refund / total_income if total_income else 0

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", 'B', 16)
    pdf.cell(200, 10, txt="经营分析报告（PDF）", ln=True, align='C')
    pdf.ln(10)

    pdf.set_font("Arial", '', 12)
    pdf.cell(200, 10, txt=f"总收入：¥{total_income:.2f}", ln=True)
    pdf.cell(200, 10, txt=f"实际收入：¥{actual:.2f}", ln=True)
    pdf.cell(200, 10, txt=f"退款金额：¥{refund:.2f}（退款率：{refund_rate:.2%}）", ln=True)

    pdf.ln(10)
    pdf.set_font("Arial", 'B', 12)
    pdf.cell(200, 10, txt="建议与结论：", ln=True)
    pdf.set_font("Arial", '', 12)
    if refund_rate > 0.15:
        pdf.multi_cell(0, 10, "退款率偏高，建议优化套餐说明和服务流程")
    else:
        pdf.multi_cell(0, 10, "退款率处于正常范围")

    out_path = os.path.join(REPORT_FOLDER, f"report_{uuid.uuid4().hex}.pdf")
    pdf.output(out_path)
    return out_path

@app.route('/logout')
def logout():
    session.clear()
    return redirect('/')

# ✅ 正确启动（适配 Railway）
if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port)
